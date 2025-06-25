import argparse
import os
import re
import tempfile
from itertools import combinations

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FRACTION_PHRASES = r"(?:\s*out\s+of\s*|\s*/\s*|\s*of\s*|\s*in\s*)"
FRACTION_PATTERN = re.compile(rf"(\d+(?:\.\d+)?)\s*{FRACTION_PHRASES}\s*(\d+(?:\.\d+)?)", re.IGNORECASE)

def extract_percentage_from_string(value):
    try:
        if isinstance(value, str):
            match = FRACTION_PATTERN.search(value.lower())
            if match:
                numerator = float(match.group(1))
                denominator = float(match.group(2))
                if denominator != 0:
                    return round((numerator / denominator) * 100, 2)
        return value
    except Exception:
        return value

def read_excel_with_dynamic_header(file_path, sheet_name=0, id_column="id"):
    try:
        raw_data = pd.read_excel(file_path, sheet_name=sheet_name, header=None, keep_default_na=False)
        header_row_index = None

        for idx, row in raw_data.iterrows():
            if any(str(cell).strip().lower() == id_column.lower() for cell in row.values):
                header_row_index = idx
                break

        if header_row_index is None:
            raise ValueError(f"Could not find the column '{id_column}' in the file.")

        df = pd.read_excel(file_path, sheet_name=sheet_name, header=header_row_index, keep_default_na=False)
        for col in df.select_dtypes(include='object').columns:
            df[col] = df[col].apply(extract_percentage_from_string)

        return df

    except Exception as e:
        print(f"Failed to read Excel file: {e}")
        return None

def get_categorical_columns(df, max_unique=5, min_unique=None, dropna=True):
    if min_unique is None:
        return [col for col in df.columns if df[col].nunique(dropna=dropna) <= max_unique]
    return [col for col in df.columns if min_unique <= df[col].nunique(dropna=dropna) <= max_unique]

def replace_minor_non_numeric_with_nan(df, threshold=4):
    df = df.copy()
    categorical_cols = get_categorical_columns(df)  # Default max_unique=5
    for col in df.columns:
        if col not in categorical_cols:
            try:
                numeric_coerced = pd.to_numeric(df[col], errors='coerce')
                non_numeric_mask = numeric_coerced.isna() & df[col].notna()
                non_numeric_uniques = df[col][non_numeric_mask].unique()
                if 0 < len(non_numeric_uniques) < threshold:
                    print(f"Column '{col}': replacing values {non_numeric_uniques} with NaN")
                    df[col] = numeric_coerced
            except Exception as e:
                print(f"Skipping column '{col}' due to error: {e}")
    return df



def save_plot(fig, title, tmpdir, dpi=200):
    fname = os.path.join(tmpdir, f"{re.sub(r'[^a-zA-Z0-9]', '_', title)}.png")
    fig.savefig(fname, dpi=dpi, bbox_inches='tight')
    plt.close(fig)
    return fname

def display_categorical_statistics(df, columns, dropna=True):
    stats = []
    for col in columns:
        value_counts = df[col].value_counts(dropna=dropna)
        mode = value_counts.idxmax()
        prob_dist = value_counts / value_counts.sum()
        expected = ', '.join([f"{cat}*{prob:.2f}" for cat, prob in prob_dist.items()])
        stat = {
            'column': col,
            'mode': mode,
            'expected': expected,
            'prob_dist': prob_dist.to_dict()
        }
        stats.append(stat)
    return stats

def plot_categorical_pie_charts(df, columns, tmpdir, dropna=True):
    results = []
    for col in columns:
        fig, ax = plt.subplots(figsize=(8,6))
        df[col].value_counts(normalize=True, dropna=dropna).plot(
            kind='pie', autopct='%1.2f%%', ax=ax
        )
        ax.set_title(f"{col}")
        ax.set_ylabel("")
        img_path = save_plot(fig, f"pie_{col}", tmpdir)
        caption = f"Distribution of {col}"
        results.append((img_path, caption))
    return results

def plot_contingency_heatmaps(df, columns, tmpdir, dropna=True, as_percent=True):
    images = []
    for col1, col2 in combinations(columns, 2):
        counts = pd.crosstab(df[col1], df[col2], dropna=dropna)
        if as_percent:
            contingency = counts / counts.sum().sum() * 100
            fmt = ".2f"
        else:
            contingency = counts
            fmt = "d"
        fig, ax = plt.subplots(figsize=(8,6))
        sns.heatmap(contingency, annot=True, fmt=fmt, cmap="YlGnBu", ax=ax)
        ax.set_title(f"{col1} vs {col2}" + (" (Percent)" if as_percent else ""))
        ax.set_xlabel(col2)
        ax.set_ylabel(col1)
        images.append(save_plot(fig, f"contingency_{col1}_{col2}", tmpdir))
    return images

def plot_high_cardinality_categoricals(df, tmpdir, min_unique=6, max_unique=30, max_label_length=25, dropna=True):
    images = []
    categorical_cols = get_categorical_columns(df, min_unique=min_unique, max_unique=max_unique, dropna=dropna)
    for col in categorical_cols:
        temp = df[col].astype(str).apply(lambda x: x[:max_label_length] + '…' if len(x) > max_label_length else x)
        fig, ax = plt.subplots(figsize=(8,6))
        sns.countplot(y=temp,hue=temp, order=temp.value_counts().index, palette="Set2", ax=ax)
        ax.set_title(f"{col}")
        for label in ax.get_xticklabels():
            label.set_rotation(45)
            label.set_ha('right')
        images.append(save_plot(fig, f"high_cardinality_{col}", tmpdir))
    return images

def plot_boxplot(data, col, mean_val, median_val, use_log, tmpdir):
    fig, ax = plt.subplots(figsize=(8,6))
    sns.boxplot(x=data, color='lightgreen', ax=ax)
    ax.axvline(mean_val, color='blue', linestyle='--', label='Mean')
    ax.axvline(median_val, color='red', linestyle='--', label='Median')
    ax.set_title(f'{col} {"(log scale)" if use_log else ""}')
    ax.set_xlabel(col)
    ax.legend()
    fname = save_plot(fig, f"box_{col}", tmpdir)
    return fname

def plot_histogram_with_kde(data, col, mean_val, median_val, std_val, use_log, tmpdir):
    fig, ax = plt.subplots(figsize=(8,6))
    sns.histplot(data, kde=True, color='cornflowerblue', bins=30, ax=ax)
    ax.axvline(mean_val, color='blue', linestyle='--', label='Mean')
    ax.axvline(median_val, color='red', linestyle='--', label='Median')
    ax.axvline(mean_val + std_val, color='purple', linestyle=':', label='±1 Std Dev')
    ax.axvline(mean_val - std_val, color='purple', linestyle=':')
    ax.set_title(f'Histogram of {col} {"(log scale)" if use_log else ""}')
    ax.set_xlabel(col)
    ax.set_ylabel('Frequency')
    ax.legend()
    fname = save_plot(fig, f"hist_{col}", tmpdir)
    return fname

def plot_hexbin(x_data, y_data, x, y, use_log_x, use_log_y, tmpdir):
    fig, ax = plt.subplots(figsize=(8,6))
    ax.hexbin(x_data, y_data, gridsize=30, cmap='YlGnBu', bins='log')
    cb = fig.colorbar(ax.collections[0], ax=ax, label='log(count)')
    ax.set_xlabel(x + (" (log)" if use_log_x else ""))
    ax.set_ylabel(y + (" (log)" if use_log_y else ""))
    ax.set_title(f"{x} vs {y}" + (" (log transformed)" if use_log_x or use_log_y else ""))
    fname = save_plot(fig, f"hexbin_{x}_{y}", tmpdir)
    return fname

def plot_kde_contour(x_data, y_data, x, y, use_log_x, use_log_y, tmpdir):
    fig, ax = plt.subplots(figsize=(8,6))
    sns.kdeplot(x=x_data, y=y_data, ax=ax)
    ax.set_xlabel(x + (" (log)" if use_log_x else ""))
    ax.set_ylabel(y + (" (log)" if use_log_y else ""))
    ax.set_title(f"{x} vs {y}" + (" (log scale applied)" if use_log_x or use_log_y else ""))
    fname = save_plot(fig, f"kde_{x}_{y}", tmpdir)
    return fname




def add_flexible_slide(prs, title="Title Placeholder", caption="Caption Placeholder", image_path=None):
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # ─────────────────────────────
    # Title (Top - Full Width)
    # ─────────────────────────────
    margin = Inches(1)
    title_height = Inches(1)
    title_shape = slide.shapes.add_textbox(margin, Inches(0.2), prs.slide_width - 2 * margin, title_height)
    title_tf = title_shape.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # ─────────────────────────────
    # Content Area (Below Title)
    # ─────────────────────────────
    content_top = Inches(1.4)
    content_height = Inches(5.0)

    # Caption Box (Left Side)
    caption_width = Inches(4)
    caption_shape = slide.shapes.add_textbox(margin, content_top, caption_width, content_height)
    caption_tf = caption_shape.text_frame
    caption_tf.word_wrap = True
    for idx, line in enumerate(caption.split('\n')):
        p = caption_tf.add_paragraph() if idx > 0 else caption_tf.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.LEFT

    # Image (Right Side)
    if image_path:
        img_left = margin + caption_width + Inches(0.5)  # spacing between caption and image
        img_width = prs.slide_width - img_left - margin
        img_height = content_height
        slide.shapes.add_picture(image_path, img_left, content_top, width=img_width, height=img_height)

    return slide


def generate_ppt_from_excel(input_file, output_file):
    with tempfile.TemporaryDirectory() as tmpdir:
        df = read_excel_with_dynamic_header(input_file)
        if df is None:
            print("Failed to read Excel file.")
            return
        df = replace_minor_non_numeric_with_nan(df)

        prs = Presentation("Presentation 4.pptx")

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Website Analysis Report"
        # title_slide.placeholders[1].text = f"Input File: {os.path.basename(input_file)}"

        # Categorical analysis
        cat_cols = get_categorical_columns(df)
        cat_stats = display_categorical_statistics(df, cat_cols)
        pie_imgs = plot_categorical_pie_charts(df, cat_cols, tmpdir)
        for (img_path, caption), stat in zip(pie_imgs, cat_stats):
            # bullet list as caption
            bullets = [
                f"Mode: {stat['mode']}",
                f"Expected: {stat['expected']}",
            ] + [f"{cat}: {prob:.2f}" for cat, prob in stat['prob_dist'].items()]
            add_flexible_slide(prs, f"Distribution: {stat['column']}", "\n".join(bullets), img_path)

        cont_imgs = plot_contingency_heatmaps(df, cat_cols, tmpdir)
        for img in cont_imgs:
            add_flexible_slide(prs, "Contingency Heatmap", "Joint probabilities between variables.", img)

        high_card_imgs = plot_high_cardinality_categoricals(df, tmpdir)
        for img in high_card_imgs:
            add_flexible_slide(prs, "High Cardinality Categorical", "Bar plot for high-cardinality features.", img)

        num_cols = df.select_dtypes(include=['number']).columns
        if len(num_cols) == 0:
            add_flexible_slide(prs, "Numerical Analysis", "No numerical columns found in the DataFrame.", None)
        else:
            num_cols = [col for col in num_cols if df[col].nunique() < len(df)]
            for col in num_cols:
                plot_data = df[col].copy()
                use_log = plot_data.max() > 1e6
                if use_log:
                    plot_data = np.log1p(plot_data)
                mean_val = plot_data.mean()
                median_val = plot_data.median()
                std_val = plot_data.std()
                stats_lines = [
                    f"Mean: {mean_val:.2f}",
                    f"Median: {median_val:.2f}",
                    f"Std Dev: {std_val:.2f}"
                ]
                img = plot_boxplot(plot_data, col, mean_val, median_val, use_log, tmpdir)
                add_flexible_slide(prs, f"Boxplot of {col}", "\n".join(stats_lines), img)
                img = plot_histogram_with_kde(plot_data, col, mean_val, median_val, std_val, use_log, tmpdir)
                add_flexible_slide(prs, f"Histogram of {col}", "\n".join(stats_lines), img)
            # Pairwise plots
            for i, x in enumerate(num_cols):
                for y in num_cols[i+1:]:
                    x_data = df[x].copy()
                    y_data = df[y].copy()
                    use_log_x = x_data.max() > 1e6
                    use_log_y = y_data.max() > 1e6
                    if use_log_x:
                        x_data = np.log1p(x_data)
                    if use_log_y:
                        y_data = np.log1p(y_data)
                    img = plot_hexbin(x_data, y_data, x, y, use_log_x, use_log_y, tmpdir)
                    add_flexible_slide(prs, f"Hexbin: {x} vs {y}", f"Joint distribution of {x} and {y}.", img)
                    img = plot_kde_contour(x_data, y_data, x, y, use_log_x, use_log_y, tmpdir)
                    add_flexible_slide(prs, f"KDE Contour: {x} vs {y}", f"KDE Contour of {x} and {y}.", img)

        # Final slide
        add_flexible_slide(prs, "End of Report", "This concludes the automated analysis.", None)
        for slide in prs.slides:
            for shape in list(slide.shapes):  # list() prevents mutation issues
                if shape.is_placeholder:
                    slide.shapes._spTree.remove(shape._element)
        prs.save(output_file)
        print(f"Presentation generated at: {output_file}")

# --- CLI Entry Point ---

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Excel Data Analysis to PPTX")
    parser.add_argument("input_excel", help="Path to input Excel file")
    parser.add_argument("output_pptx", help="Path to output PPTX file")
    args = parser.parse_args()
    generate_ppt_from_excel(args.input_excel, args.output_pptx)    