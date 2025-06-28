
import argparse
import asyncio
import io
import os
import random
import re
import sys
from typing import Any, List, Optional, Tuple

import matplotlib.cm as cm
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ========== Data Loading and Cleaning ==========

NA_VALUES = ['N/A', 'n/a', 'NA', 'N.A.', 'not applicable', '-']
INACTIVE_SITE_INDICATORS = ['N/A - Dead Links / Redirects', 'Dead Links', 'Redirects']

def find_header_row(file_path, num_rows_to_check=50):
    """Detects the header row in an Excel file."""
    try:
        temp_df = pd.read_excel(file_path, header=None, nrows=num_rows_to_check)
        best_header_row_index = 0
        max_score = -1
        for i in range(len(temp_df)):
            potential_header_series = temp_df.iloc[i]
            numeric_check = pd.to_numeric(potential_header_series, errors='coerce')
            non_numeric_count = numeric_check.isna().sum()
            row_as_strings = potential_header_series.astype(str)
            unique_non_empty_strings_count = row_as_strings[row_as_strings.str.strip() != ''].nunique()
            current_score = (unique_non_empty_strings_count * 3) + non_numeric_count
            if current_score > max_score:
                max_score = current_score
                best_header_row_index = i
        print(f"[INFO] Detected header row (0-indexed): {best_header_row_index}")
        return best_header_row_index
    except Exception as e:
        print(f"[WARN] Error during header detection: {e}")
        return 0

def load_and_clean_data(file_path: str) -> pd.DataFrame:
    """Loads and cleans the data from Excel."""
    header_row_index = find_header_row(file_path)
    df = pd.read_excel(file_path, header=header_row_index)
    df.columns = df.columns.str.strip()
    df_cleaned = df.copy()

    for col in df_cleaned.select_dtypes(include='object').columns:
        df_cleaned[col] = df_cleaned[col].astype(str)
        df_cleaned[col] = df_cleaned[col].replace(NA_VALUES, np.nan)
        df_cleaned[col] = df_cleaned[col].replace(r'^\s*$', np.nan, regex=True)
    print("[INFO] Replaced common 'N/A' and empty strings with NaN.")
    traffic_columns = [col for col in df_cleaned.columns if "traffic" in col.lower()]
    if traffic_columns:
        active_mask = pd.Series([True] * len(df_cleaned), index=df_cleaned.index)
        for col in traffic_columns:
            col_inactive_mask = df_cleaned[col].astype(str).isin(INACTIVE_SITE_INDICATORS)
            active_mask = active_mask & (~col_inactive_mask)
        active_sites = df_cleaned[active_mask].copy()
        print(f"[INFO] Filtered out inactive sites using columns: {', '.join(traffic_columns)}.")
    else:
        print("[WARN] No columns with 'traffic' found. Skipping inactive site filtering.")
        active_sites = df_cleaned.copy()
    return active_sites

def parse_score_advanced(score_str):
    if pd.isna(score_str) or str(score_str).strip() == '' or str(score_str).strip().lower() in [s.lower() for s in NA_VALUES]:
        return np.nan
    s = str(score_str).strip().lower()
    # "X out of Y"
    match_out_of = re.match(r'(\d+(?:\.\d+)?)\s*out of\s*(\d+(?:\.\d+)?)', s)
    if match_out_of:
        try:
            numerator = float(match_out_of.group(1))
            denominator = float(match_out_of.group(2))
            return (numerator / denominator) * 100.0 if denominator != 0 else np.nan
        except ValueError:
            pass
    # "N%"
    match_percent = re.match(r'(\d+(?:\.\d+)?)\s*%', s)
    if match_percent:
        try:
            return float(match_percent.group(1))
        except ValueError:
            pass

    try:
        return float(s)
    except ValueError:
        pass
    return np.nan


def identify_id_columns_by_pattern(df, patterns=None, verbose=False):
    
    if patterns is None:
        patterns = [
            r'\b(?:id|no|num|number|key|serial|code|idx)\b', 
            r'_id\b',       
            r'id_',          
            r'sl_no',       
            r'reference',  
            r'ref\b'        
        ]
    elif isinstance(patterns, str):
        patterns = [patterns] 

    potential_id_columns = set()

    for col in df.columns:
        original_col_name = col
        normalized_col_name = normalize_col_name(col)

        if verbose:
            print(f"\n--- Analyzing column: '{original_col_name}' (Normalized: '{normalized_col_name}') ---")

        for pattern in patterns:
            if re.search(pattern, normalized_col_name, re.IGNORECASE):
                potential_id_columns.add(original_col_name)
                if verbose:
                    print(f"  - Flagged by: Column name matches pattern '{pattern}'")
                break

    return list(potential_id_columns)


def clean_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """Cleans and normalizes the loaded DataFrame."""

    name_col = next((col for col in df.columns if "website" in col.lower() or "domain" in col.lower()), None)
    if name_col:
        df = df.drop_duplicates(subset=[name_col], keep='first')
        print(f"[INFO] Removed duplicates based on '{name_col}'.")

    categorical_cols = df.select_dtypes(include=['object']).columns
    for col_name in categorical_cols:
        df[col_name] = df[col_name].astype(str).str.strip()
        df[col_name] = df[col_name].replace(r'^\s*$', np.nan, regex=True)

    score_columns = [col for col in df.columns if 'score' in col.lower()]
    for col_name in score_columns:
        df[col_name] = df[col_name].astype(str).apply(parse_score_advanced)

    for col in df.columns:
        converted_col = pd.to_numeric(df[col], errors='coerce')
        if not converted_col.isnull().all() and converted_col.dtype != df[col].dtype:
            df[col] = converted_col
    print(f"[INFO] DataFrame cleaned. Final shape: {df.shape}")
    return df

# ========== Analysis Helpers ==========

def normalize_col_name(col_name):
    normalized = col_name.lower().strip()
    normalized = normalized.replace(' / ', '/').replace(' ', '_')
    return normalized

def describe_distribution(series: pd.Series) -> str:
    skew = series.skew()
    if skew > 0.5:
        return "right-skewed (most values below average)"
    elif skew < -0.5:
        return "left-skewed (most values above average)"
    else:
        return "normally distributed"

def get_risk_counts(scores: pd.Series) -> Tuple[int, int, int]:
    high = (scores < 60).sum()
    medium = ((scores >= 60) & (scores < 90)).sum()
    low = (scores >= 90).sum()
    return high, medium, low

def generate_bullet_points_for_chart(df: pd.DataFrame, col: str, chart_type: str) -> List[str]:
    """Generate contextual bullet points based on data analysis."""
    bullet_points = []
    nan_count_current_df = df[col].isnull().sum() + (df[col].astype(str).str.lower() == 'nan').sum()
    total_records_current_df = len(df)
    nan_pct_current_df = (nan_count_current_df / total_records_current_df) * 100 if total_records_current_df > 0 else 0
    valid_data_filtered = df[col].dropna()
    valid_data_filtered = valid_data_filtered[valid_data_filtered.astype(str).str.lower() != 'nan']
    if chart_type == "risk_distribution" and not valid_data_filtered.empty:
        high, medium, low = get_risk_counts(valid_data_filtered)
        total = len(valid_data_filtered)
        high_pct = (high / total) * 100 if total > 0 else 0
        medium_pct = (medium / total) * 100 if total > 0 else 0
        low_pct = (low / total) * 100 if total > 0 else 0
        mean_score = valid_data_filtered.mean()
        bullet_points = [
            f"Total {total} records analyzed with average score of {mean_score:.1f}",
            f"High risk sites represent {high_pct:.1f}% ({high}) of total records",
            f"Medium risk sites account for {medium_pct:.1f}% ({medium}) of total records",
            f"Low risk sites account for {low_pct:.1f}% ({low}) of total records",
        ]
    elif chart_type == "score_distribution" and not valid_data_filtered.empty:
        mean, median, std = valid_data_filtered.mean(), valid_data_filtered.median(), valid_data_filtered.std()
        dist_desc = describe_distribution(valid_data_filtered)
        min_val, max_val = valid_data_filtered.min(), valid_data_filtered.max()
        bullet_points = [
            f"Scores range from {min_val:.1f} to {max_val:.1f} with a mean of {mean:.1f}.",
            f"Standard deviation of {std:.1f} shows {'high' if std > mean * 0.3 else 'moderate'} variability.",
            f"The distribution is {dist_desc}.",
            f"Median of {median:.1f} is {'close to' if abs(mean-median)<std*0.1 else 'distinct from'} the mean."
        ]
    elif chart_type == "numerical_distribution" and not valid_data_filtered.empty:
        mean, median = valid_data_filtered.mean(), valid_data_filtered.median()
        q25, q75 = valid_data_filtered.quantile(0.25), valid_data_filtered.quantile(0.75)
        dead_links_count = nan_count_current_df
        dead_links_pct = nan_pct_current_df
        bullet_points = [
            f"The dataset has {len(valid_data_filtered)} valid entries and {dead_links_count} missing entries (about {dead_links_pct:.1f}%).",
            f"The average value is {mean:.2f}, with a median of {median:.2f}.",
            f"The 25th percentile is {q25:.2f}, and the 75th percentile is {q75:.2f}.",
            f"The middle 50% of data lies between {q25:.2f} and {q75:.2f}."
        ]
    elif chart_type in {"categorical_pie", "categorical_bar"}:
        if nan_count_current_df > 0:
            bullet_points.append(
                f"{nan_pct_current_df:.1f}% of entries ({nan_count_current_df} records) are missing or 'nan'."
            )
        if not valid_data_filtered.empty:
            value_counts = valid_data_filtered.value_counts()
            total_count_filtered = len(valid_data_filtered)
            if not value_counts.empty:
                top_category, top_count = value_counts.index[0], value_counts.iloc[0]
                top_pct = (top_count / total_count_filtered) * 100
                bullet_points.append(f"The dataset contains {len(value_counts)} distinct categories across {total_count_filtered} valid records (excluding 'nan').")
                bullet_points.append(f"The dominant category '{top_category}' has {top_count} occurrences ({top_pct:.1f}%).")
                if chart_type == "categorical_pie":
                    bullet_points.append(f"The distribution is {'relatively uniform' if top_pct < 40 else 'skewed towards a few dominant categories'}.")
                elif chart_type == "categorical_bar":
                    top_10_count_filtered = value_counts.head(10).sum()
                    top_10_make_up_pct = (top_10_count_filtered / total_count_filtered) * 100
                    bullet_points.append(f"The top 10 categories (excluding 'nan') make up {top_10_make_up_pct:.1f}% of all valid entries.")
                bullet_points.append(f"The categories are {'spread out across many values' if len(value_counts) > total_count_filtered * 0.5 else 'mostly focused on a few values'}.")
            else:
                bullet_points.append(f"No valid (non-NaN) categorical data found for {col}.")
        else:
            bullet_points.append(f"No valid (non-NaN) categorical data found for {col}.")
    return bullet_points[:4]

def generate_hexbin_bullet_points(df: pd.DataFrame, x_col: str, y_col: str) -> List[str]:
    corr = df[x_col].corr(df[y_col])
    return [
        "Darker hexagons indicate denser data regions, while lighter ones show sparser observations.",
        f"A {'positive' if corr > 0 else 'negative' if corr < 0 else 'no'} linear trend exists (Pearson correlation coefficient: {corr:.2f}).",
        f"{x_col} ranges from {df[x_col].min():.1f} to {df[x_col].max():.1f}, and {y_col} ranges from {df[y_col].min():.1f} to {df[y_col].max():.1f}."
    ]

def add_conclusion(df: pd.DataFrame) -> List[str]:
    volume_cols = [col for col in df.columns if "volume" in col.lower()]
    str_volume_cols = [col for col in volume_cols if df[col].apply(lambda x: isinstance(x, str)).any()]
    score_cols = [col for col in df.columns if "score" in col.lower()]
    name_cols = [col for col in df.columns if any(keyword in col.lower() for keyword in ["website", "domain", "site", "link"])]
    matching_rows = pd.DataFrame()
    for vol_col in str_volume_cols:
        for score_col in score_cols:
            if pd.api.types.is_numeric_dtype(df[score_col]):
                mask = df[vol_col].str.contains("high", case=False, na=False) & (df[score_col] < 50)
                filtered = df.loc[mask, name_cols]
                matching_rows = pd.concat([matching_rows, filtered], ignore_index=True)
    result = matching_rows.drop_duplicates().head(10)
    if not result.empty and name_cols:
        first_name_col = result.columns[0]
        top_sites = result[first_name_col].dropna().astype(str).unique()[:9]
        if len(top_sites) > 0:
            if len(top_sites) > 1:
                site_list = ', '.join(top_sites[:-1]) + f', and {top_sites[-1]}'
            else:
                site_list = top_sites[0]
            score_col_name = score_cols[0].lower() if score_cols else "score"
            volume_col_name = str_volume_cols[0].lower() if str_volume_cols else "volume"
            return [
                f"The current dataset includes the following key columns: {', '.join(df.columns)}",
                f"The sites {site_list} are considered high-risk given their low {score_col_name} despite having high {volume_col_name}."
            ]
    return [f"The current dataset includes the following key columns: {', '.join(df.columns)}"]

# ========== Visualization Functions ==========

def plot_numerical_histogram(ax, series: pd.Series, col: str, clean_title: str):
    mean_val, median_val, std_val = series.mean(), series.median(), series.std()
    sns.histplot(series, kde=True, color='cornflowerblue', bins=30, ax=ax)
    ax.axvline(mean_val, color='blue', linestyle='--', label=f'Mean: {mean_val:.2f}')
    ax.axvline(median_val, color='red', linestyle='--', label=f'Median: {median_val:.2f}')
    ax.axvline(mean_val + std_val, color='purple', linestyle=':', label=f'+-1 Std Dev: {std_val:.2f}')
    ax.axvline(mean_val - std_val, color='purple', linestyle=':')
    ax.set_title(f'Distribution of {clean_title}', fontsize=16, pad=20)
    ax.set_xlabel(col, fontsize=12)
    ax.set_ylabel('Frequency', fontsize=12)
    ax.legend()

def plot_categorical_bar(ax, value_counts: pd.Series, col: str, clean_title: str):
    max_display_label_length = 20
    sorted_index = value_counts.reset_index()
    sorted_index.columns = ['label', 'count']
    sorted_index['label'] = sorted_index['label'].astype(str)
    sorted_index["label_length"] = sorted_index["label"].apply(len)
    sorted_index["label_lower"] = sorted_index["label"].str.lower()
    sorted_index = sorted_index.sort_values(
        by=["count", "label_length", "label_lower"],
        ascending=[False, True, True]
    ).drop(columns="label_lower")
    top_10_df = sorted_index.head(10).copy()
    top_10_df["short_label"] = top_10_df["label"].apply(
        lambda x: x if len(x) <= max_display_label_length else x[:max_display_label_length] + "…"
    )
    sns.barplot(
        x=top_10_df["count"],
        y=top_10_df["short_label"],
        hue=top_10_df["short_label"],
        palette="Set2",
        orient="h",
        ax=ax,
        legend=False
    )
    ax.set_title(f'Top 10 Most Frequent Values in {clean_title} (Excluding NaN)', fontsize=16, pad=20)
    ax.set_xlabel('Count', fontsize=12)
    ax.set_ylabel(col, fontsize=12)

def visualize_column_summary(active_sites_df: pd.DataFrame) -> Tuple[List[Tuple[Any, str, List[str]]], List[str]]:
    sns.set_style("whitegrid")
    color_palette = ['#9370DB', '#FF6347', '#FFB300', '#32CD32', '#27AEEF']
    
    raw_cols_to_exclude = identify_id_columns_by_pattern(active_sites_df, verbose=True)
    normalized_cols_to_exclude_set = {normalize_col_name(col) for col in raw_cols_to_exclude}
    
    score_keyword = 'score'
    cols_for_viz = [col for col in active_sites_df.columns if normalize_col_name(col) not in normalized_cols_to_exclude_set]
    chart_data = []
    for col in cols_for_viz:
        clean_title = col.replace('_', ' ').replace('-', ' ').title()
        # Numerical Columns
        if pd.api.types.is_numeric_dtype(active_sites_df[col]):
            valid_numerical_data = active_sites_df[col].dropna()
            valid_numerical_data = valid_numerical_data[valid_numerical_data.astype(str).str.lower() != 'nan']
            if score_keyword in normalize_col_name(col):
                if not valid_numerical_data.empty:
                    # Pie chart: risk distribution for score columns
                    fig1, ax1 = plt.subplots(figsize=(10, 6))
                    high, medium, low = get_risk_counts(valid_numerical_data)
                    risk_counts = [high, medium, low]
                    risk_labels = ['High Risk (< 60)', 'Medium Risk (60-89)', 'Low Risk (>= 90)']
                    colors_for_risk_pie = random.sample(color_palette, k=min(len(risk_labels), len(color_palette)))
                    explode_values = [0.05 if i == 0 and high > 0 else 0 for i in range(len(risk_labels))]
                    ax1.pie(
                        risk_counts, labels=risk_labels, autopct='%1.1f%%', startangle=140,
                        colors=colors_for_risk_pie, wedgeprops={'edgecolor': 'white'}, shadow=True, explode=explode_values
                    )
                    ax1.set_title(f'{clean_title} Risk Distribution', fontsize=16, pad=20)
                    ax1.set_ylabel('')
                    fig1.tight_layout()
                    chart_data.append((fig1, f"{clean_title} Risk Distribution",
                        generate_bullet_points_for_chart(active_sites_df, col, "risk_distribution")))
                    plt.close(fig1)
                    # Histogram: score distribution
                    fig2, ax2 = plt.subplots(figsize=(10, 6))
                    plot_numerical_histogram(ax2, valid_numerical_data, col, clean_title)
                    fig2.tight_layout()
                    chart_data.append((fig2, f"{clean_title} Distribution",
                        generate_bullet_points_for_chart(active_sites_df, col, "score_distribution")))
                    plt.close(fig2)
            else: # General numerical column
                if not valid_numerical_data.empty:
                    fig, ax = plt.subplots(figsize=(10, 6))
                    plot_numerical_histogram(ax, valid_numerical_data, col, clean_title)
                    fig.tight_layout()
                    chart_data.append((fig, f"{clean_title} Distribution",
                        generate_bullet_points_for_chart(active_sites_df, col, "numerical_distribution")))
                    plt.close(fig)
        # Categorical Columns
        else:
            non_null_series_filtered_for_plotting = active_sites_df[col].dropna()
            non_null_series_filtered_for_plotting = non_null_series_filtered_for_plotting[non_null_series_filtered_for_plotting.astype(str).str.lower() != 'nan']
            unique_count = non_null_series_filtered_for_plotting.nunique()
            if unique_count == 0 or unique_count == len(non_null_series_filtered_for_plotting):
                nan_count_for_bullet = active_sites_df[col].isnull().sum() + (active_sites_df[col].astype(str).str.lower() == 'nan').sum()
                if nan_count_for_bullet == len(active_sites_df):
                    chart_data.append((None, f"{clean_title} Distribution (No valid data)",
                        generate_bullet_points_for_chart(active_sites_df, col, "categorical_bar")))
                continue
            if unique_count <= 5:
                fig, ax = plt.subplots(figsize=(10, 6))
                value_counts = non_null_series_filtered_for_plotting.value_counts()
                explode_values = [0.05] * len(value_counts)
                colors_for_categorical_pie = random.sample(color_palette, k=min(len(value_counts), len(color_palette)))
                ax.pie(
                    value_counts, labels=value_counts.index.astype(str),
                    autopct='%1.1f%%', startangle=140,
                    colors=colors_for_categorical_pie, wedgeprops={'edgecolor': 'white'}, shadow=True, explode=explode_values
                )
                ax.set_title(f'Distribution of {clean_title} (excluding nan values)', fontsize=16, pad=20)
                ax.set_ylabel('')
                fig.tight_layout()
                chart_data.append((fig, f"{clean_title} Distribution",
                    generate_bullet_points_for_chart(active_sites_df, col, "categorical_pie")))
                plt.close(fig)
            else:
                fig, ax = plt.subplots(figsize=(10, 6))
                value_counts_for_bar_plot = non_null_series_filtered_for_plotting.value_counts()
                plot_categorical_bar(ax, value_counts_for_bar_plot, col, clean_title)
                fig.tight_layout()
                chart_data.append((fig, f"Top 10 {clean_title} Values",
                    generate_bullet_points_for_chart(active_sites_df, col, "categorical_bar")))
                plt.close(fig)
    # Hexbin plot (traffic vs score)
    traffic_col, score_col = None, None
    for col in active_sites_df.columns:
        norm_col = normalize_col_name(col)
        if pd.api.types.is_numeric_dtype(active_sites_df[col]):
            if 'traffic' in norm_col:
                traffic_col = col
            elif 'score' in norm_col:
                score_col = col
    if traffic_col and score_col:
        df_hex = active_sites_df[[traffic_col, score_col]].dropna()
        df_hex = df_hex[df_hex[traffic_col].astype(str).str.lower() != 'nan']
        df_hex = df_hex[df_hex[score_col].astype(str).str.lower() != 'nan']
        if not df_hex.empty:
            fig, ax = plt.subplots(figsize=(10, 6))
            hb = ax.hexbin(df_hex[traffic_col], df_hex[score_col], gridsize=30, cmap='viridis_r', mincnt=1)
            cb = fig.colorbar(hb, ax=ax)
            cb.set_label('Count')
            ax.set_xlabel(traffic_col)
            ax.set_ylabel(score_col)
            ax.set_title(f'{traffic_col} vs {score_col} Hexbin Plot', fontsize=16, pad=20)
            fig.tight_layout()
            chart_title = f"{traffic_col.replace('_',' ').title()} vs {score_col.replace('_',' ').title()} Relationship"
            chart_data.append((fig, chart_title, generate_hexbin_bullet_points(df_hex, traffic_col, score_col)))
            plt.close(fig)
   
    conclusion_bullets = add_conclusion(active_sites_df)
    return chart_data, conclusion_bullets

# ========== PPTX Creation Functions ==========

TITLE_FONT_NAME = 'Times New Roman'
TITLE_FONT_COLOR = RGBColor(112, 48, 160)
TITLE_FONT_SIZE = Pt(32)
BULLET_FONT_NAME = 'Times New Roman'
BULLET_FONT_SIZE = Pt(20)
BULLET_SPACE_AFTER = Pt(12)
TOP_MARGIN_RATIO = 0.06
BOTTOM_MARGIN_RATIO = 0.08
LEFT_MARGIN_RATIO = 0.06
RIGHT_MARGIN_RATIO = 0.06
TITLE_HEIGHT_RATIO = 0.20
CONTENT_HEIGHT_RATIO = 1 - TITLE_HEIGHT_RATIO
TITLE_WIDTH_RATIO = 0.70
TEXT_WIDTH_RATIO = 0.30
IMAGE_WIDTH_RATIO = 0.70
GAP_BETWEEN_TEXT_AND_IMAGE_RATIO = 0.02

def bolden_values_paragraph(p, text):
    """Add bullet point with bolded numbers/percentages in the paragraph p."""
    pattern = re.compile(r"(\d[\d,\.]*%?|\([\d,\.]+\)|\b[\w\-]+(?:\.[\w\-]+)+\b)")
    last = 0
    for match in pattern.finditer(text):
        if match.start() > last:
            run = p.add_run()
            run.text = text[last:match.start()]
            run.font.bold = False
            run.font.size = BULLET_FONT_SIZE
            run.font.name = BULLET_FONT_NAME
        run = p.add_run()
        run.text = match.group(0)
        run.font.bold = True
        run.font.size = BULLET_FONT_SIZE
        run.font.name = BULLET_FONT_NAME
        last = match.end()
    if last < len(text):
        run = p.add_run()
        run.text = text[last:]
        run.font.bold = False
        run.font.size = BULLET_FONT_SIZE
        run.font.name = BULLET_FONT_NAME

def add_custom_chart_slide(prs: Presentation, chart_fig, chart_title: str, bullet_points: List[str]):
    # Layout calculations
    slide_width, slide_height = prs.slide_width, prs.slide_height
    usable_width = slide_width * (1 - LEFT_MARGIN_RATIO - RIGHT_MARGIN_RATIO)
    usable_height = slide_height * (1 - TOP_MARGIN_RATIO - BOTTOM_MARGIN_RATIO)
    usable_left = slide_width * LEFT_MARGIN_RATIO
    usable_top = slide_height * TOP_MARGIN_RATIO
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Title
    title_width = slide_width * TITLE_WIDTH_RATIO
    title_left = (slide_width - title_width) / 2
    title_top = usable_top
    title_height = usable_height * TITLE_HEIGHT_RATIO
    title_box = slide.shapes.add_textbox(
        Emu(title_left), Emu(title_top), Emu(title_width), Emu(title_height)
    )
    tf = title_box.text_frame
    tf.text = chart_title
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.runs[0]
    para.alignment = PP_ALIGN.CENTER
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.name = TITLE_FONT_NAME
    run.font.color.rgb = TITLE_FONT_COLOR
    # Bullets
    content_top = usable_top + usable_height * TITLE_HEIGHT_RATIO
    content_height = usable_height * CONTENT_HEIGHT_RATIO
    gap_width = usable_width * GAP_BETWEEN_TEXT_AND_IMAGE_RATIO
    text_left = Emu(usable_left)
    text_top = Emu(content_top)
    text_width = Emu(usable_width * TEXT_WIDTH_RATIO - gap_width / 2)
    text_height = Emu(content_height)
    bullet_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf_bullets = bullet_box.text_frame
    tf_bullets.word_wrap = True
    for idx, txt in enumerate(bullet_points):
        p = tf_bullets.paragraphs[0] if idx == 0 else tf_bullets.add_paragraph()
        p.level = 0
        p.space_after = BULLET_SPACE_AFTER
        p.alignment = PP_ALIGN.LEFT
        bolden_values_paragraph(p, u"\u2022 " + txt)
    # Image
    image_left = Emu(usable_left + usable_width * TEXT_WIDTH_RATIO + gap_width / 2)
    image_top = Emu(content_top)
    image_width = Emu(usable_width * IMAGE_WIDTH_RATIO - gap_width / 2)
    image_height = Emu(content_height)
    img_buffer = io.BytesIO()
    chart_fig.savefig(img_buffer, format='png', dpi=200, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close(chart_fig)
    img = Image.open(img_buffer)
    img_width, img_height = img.size
    max_width_px, max_height_px = int(image_width / 9525), int(image_height / 9525)
    aspect = img_height / img_width
    if max_height_px / aspect <= max_width_px:
        final_height_px = max_height_px
        final_width_px = int(final_height_px / aspect)
    else:
        final_width_px = max_width_px
        final_height_px = int(final_width_px * aspect)
    img_top_offset = int((max_height_px - final_height_px) / 2)
    img_left_offset = int((max_width_px - final_width_px) / 2)
    slide.shapes.add_picture(
        img_buffer,
        image_left + Emu(img_left_offset * 9525),
        image_top + Emu(img_top_offset * 9525),
        Emu(final_width_px * 9525),
        Emu(final_height_px * 9525)
    )

async def build_presentation_with_charts(
    template_path: str,
    chart_figures_and_titles: List[Tuple[Any, str, List[str]]],
    output_path: str,
    insight_points: Optional[List[str]] = None
) -> None:
    prs = Presentation(template_path)
    for idx in range(len(prs.slides) - 1, 0, -1):
        rId = prs.slides._sldIdLst[idx].rId
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])
        prs.part.drop_rel(rId)
    for fig_obj, chart_title, chart_bullets in chart_figures_and_titles:
        add_custom_chart_slide(prs, fig_obj, chart_title, chart_bullets)
    slide_width, slide_height = prs.slide_width, prs.slide_height
    usable_width = slide_width * (1 - LEFT_MARGIN_RATIO - RIGHT_MARGIN_RATIO)
    usable_height = slide_height * (1 - TOP_MARGIN_RATIO - BOTTOM_MARGIN_RATIO)
    usable_left = slide_width * LEFT_MARGIN_RATIO
    usable_top = slide_height * TOP_MARGIN_RATIO
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_width = slide_width * TITLE_WIDTH_RATIO
    title_left = (slide_width - title_width) / 2
    title_top = usable_top
    title_height = usable_height * TITLE_HEIGHT_RATIO
    title_box = slide.shapes.add_textbox(
        Emu(title_left), Emu(title_top), Emu(title_width), Emu(title_height)
    )
    tf = title_box.text_frame
    tf.text = "SUMMARY INSIGHTS"
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.name = TITLE_FONT_NAME
    run.font.color.rgb = TITLE_FONT_COLOR
    content_top = usable_top + usable_height * TITLE_HEIGHT_RATIO
    content_height = usable_height * CONTENT_HEIGHT_RATIO
    insight_left = Emu(usable_left + usable_width * 0.15)
    insight_top = Emu(content_top)
    insight_width = Emu(usable_width * 0.75)
    insight_height = Emu(content_height)
    insight_box = slide.shapes.add_textbox(insight_left, insight_top, insight_width, insight_height)
    tf_bullets = insight_box.text_frame
    tf_bullets.word_wrap = True
    for idx, txt in enumerate(insight_points or []):
        p = tf_bullets.paragraphs[0] if idx == 0 else tf_bullets.add_paragraph()
        p.level = 0
        p.space_after = BULLET_SPACE_AFTER
        p.alignment = PP_ALIGN.LEFT
        bolden_values_paragraph(p, u"\u2022 " + txt)
    # Thank you slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    thank_left = Emu(usable_left + usable_width * 0.15)
    thank_top = Emu(usable_top + usable_height * 0.35)
    thank_width = Emu(usable_width * 0.7)
    thank_height = Emu(usable_height * 0.3)
    thank_box = slide.shapes.add_textbox(thank_left, thank_top, thank_width, thank_height)
    tf = thank_box.text_frame
    tf.text = "THANK YOU"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.name = TITLE_FONT_NAME
    run.font.color.rgb = RGBColor(128, 0, 128)
    tf.word_wrap = True
    prs.save(output_path)
    print(f"[SUCCESS] Presentation created at '{output_path}'")

# ========== Main CLI Entrypoint ==========

async def main_execution(
    input_file: str,
    template_ppt_path: str,
):
    print(f"\n[START] Processing input file: {input_file}")


    active_sites = load_and_clean_data(input_file)
    active_sites = clean_dataframe(active_sites)
    chart_figures_and_titles, conclusion_points = visualize_column_summary(active_sites)
    print(f"[INFO] Chart data generated. Creating PowerPoint...")

    base, _ = os.path.splitext(input_file)
    output_file = f"{base}.pptx"

    await build_presentation_with_charts(
        template_ppt_path,
        chart_figures_and_titles,
        output_file,
        conclusion_points
    )

def is_excel_file(filename: str):
    return filename.lower().endswith(('.xlsx', '.xls'))

def is_pptx_file(filename: str):
    return filename.lower().endswith('.pptx')

def parse_args():
    parser = argparse.ArgumentParser(description="Audit Excel data and auto-generate a PowerPoint report with charts.")
    parser.add_argument('--input-file', '-i', required=True, help="Path to input Excel file (.xlsx/.xls)")
    parser.add_argument('--template', '-t', required=True, help="Path to PPTX template file (.pptx)")
    return parser.parse_args()

def run():
    args = parse_args()
    if not is_excel_file(args.input_file):
        print("[ERROR] Input file must be an Excel file (.xlsx or .xls).")
        sys.exit(1)
    if not is_pptx_file(args.template):
        print("[ERROR] Template file must be a PowerPoint file (.pptx).")
        sys.exit(1)
    if not os.path.exists(args.input_file):
        print("[ERROR] Input Excel file does not exist.")
        sys.exit(1)
    if not os.path.exists(args.template):
        print("[ERROR] Template PowerPoint file does not exist.")
        sys.exit(1)
    try:
        asyncio.run(main_execution(args.input_file, args.template))
    except Exception as e:
        print(f"[ERROR] {e}")

if __name__ == '__main__':
    run()